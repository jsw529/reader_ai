import streamlit as st
from google.cloud import texttospeech
from google.oauth2 import service_account
import tempfile
from pptx import Presentation

# client를 미리 None으로 정의 (초기화)
client = None

# 인증 파일 업로드
uploaded_service_account_file = st.file_uploader("Google Cloud 인증 JSON 파일을 업로드하세요", type=["json"])

if uploaded_service_account_file is not None:
    # 인증 파일을 저장
    with open("/mnt/data/service_account.json", "wb") as f:
        f.write(uploaded_service_account_file.getbuffer())

    # 인증 파일로 credentials 만들기
    credentials = service_account.Credentials.from_service_account_file('/mnt/data/service_account.json')

    # client 생성
    client = texttospeech.TextToSpeechClient(credentials=credentials)

    st.success("Google Cloud 인증 파일이 성공적으로 업로드되었습니다!")

# 음성 성별 선택 (남성, 여성)
voice_choice = st.radio("음성 성별 선택", ("남성", "여성"))

# 음성 설정
voice_gender = texttospeech.SsmlVoiceGender.MALE if voice_choice == "남성" else texttospeech.SsmlVoiceGender.FEMALE
voice = texttospeech.VoiceSelectionParams(
    language_code="ko-KR", ssml_gender=voice_gender
)

# PPT 파일 업로드
st.title('PPT 대본 수정 및 음성 변환 웹앱')

uploaded_ppt_file = st.file_uploader("PPT 파일 업로드", type="pptx")
if uploaded_ppt_file is not None:
    presentation = Presentation(uploaded_ppt_file)

    # 슬라이드 텍스트 추출
    slide_scripts = {}
    for i, slide in enumerate(presentation.slides):
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
        slide_scripts[i] = st.text_area(f"슬라이드 {i+1} 대본 수정", text)

    if st.button("음성 파일 생성"):
        if client is None:
            st.error("❗ 먼저 인증 파일을 업로드하세요!")
        else:
            for i, script in slide_scripts.items():
                # 텍스트를 음성으로 변환
                synthesis_input = texttospeech.SynthesisInput(text=script)
                audio_config = texttospeech.AudioConfig(
                    audio_encoding=texttospeech.AudioEncoding.MP3
                )
                response = client.synthesize_speech(
                    input=synthesis_input, voice=voice, audio_config=audio_config
                )

                # 음성 파일 생성
                with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_audio:
                    temp_audio.write(response.audio_content)
                    temp_audio_path = temp_audio.name
                    st.audio(temp_audio_path)

            st.success("🎉 음성 파일이 생성되었습니다!")
